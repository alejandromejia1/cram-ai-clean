import streamlit as st
import os
from PyPDF2 import PdfReader
from pptx import Presentation
import pytesseract
from PIL import Image
import io
import uuid
from openai import OpenAI

# File Processor Class
class FileProcessor:
    @staticmethod
    def extract_pdf_text(file):
        pdf_reader = PdfReader(file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text
    
    @staticmethod
    def extract_ppt_text(file):
        prs = Presentation(file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    
    @staticmethod
    def extract_image_text(file):
        image = Image.open(file)
        text = pytesseract.image_to_string(image)
        return text
    
    @staticmethod
    def process_file(uploaded_file):
        file_type = uploaded_file.type
        
        if file_type == "application/pdf":
            return FileProcessor.extract_pdf_text(uploaded_file)
        elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return FileProcessor.extract_ppt_text(uploaded_file)
        elif file_type.startswith('image'):
            return FileProcessor.extract_image_text(uploaded_file)
        else:
            return "Unsupported file type"

# RAG System Class
class SimpleRAG:
    def __init__(self):
        self.documents = {}
        self.current_doc_id = None
        self.client = None
        self.model_name = None
        self.conversations = {}  # Store conversations per document
        
        if 'GROQ_API_KEY' in st.secrets:
            api_key = st.secrets['GROQ_API_KEY']
            
            try:
                self.client = OpenAI(
                    api_key=api_key,
                    base_url="https://api.groq.com/openai/v1"
                )
                models = self.client.models.list()
                model_list = [model.id for model in models]
                
                if "llama-3.1-8b-instant" in model_list:
                    self.model_name = "llama-3.1-8b-instant"
                else:
                    self.model_name = model_list[0] if model_list else None
                    
                st.success("✅ System ready!")
                
            except Exception as e:
                st.error(f"API connection failed: {str(e)}")
                self.client = None
                self.model_name = None
        else:
            st.error("GROQ_API_KEY not found in secrets")
            self.client = None
            self.model_name = None
    
    def is_ready(self):
        return self.client is not None and self.model_name is not None
    
    def add_document(self, text, doc_id, filename):
        if text and text != "Unsupported file type":
            self.documents[doc_id] = {
                'text': text,
                'filename': filename,
                'processed': True
            }
            # Initialize conversation history for this document
            if doc_id not in self.conversations:
                self.conversations[doc_id] = []
            
            if self.current_doc_id is None:
                self.current_doc_id = doc_id
    
    def switch_document(self, doc_id):
        if doc_id in self.documents:
            self.current_doc_id = doc_id
    
    def delete_document(self, doc_id):
        if doc_id in self.documents:
            del self.documents[doc_id]
            if doc_id in self.conversations:
                del self.conversations[doc_id]
            if self.current_doc_id == doc_id:
                self.current_doc_id = list(self.documents.keys())[0] if self.documents else None
    
    def get_current_document(self):
        if self.current_doc_id and self.current_doc_id in self.documents:
            return self.documents[self.current_doc_id]['text']
        return None
    
    def get_conversation_history(self):
        if self.current_doc_id and self.current_doc_id in self.conversations:
            return self.conversations[self.current_doc_id]
        return []
    
    def add_to_conversation(self, question, answer):
        if self.current_doc_id:
            if self.current_doc_id not in self.conversations:
                self.conversations[self.current_doc_id] = []
            self.conversations[self.current_doc_id].append({
                'question': question,
                'answer': answer,
                'timestamp': st.session_state.get('conversation_count', 0)
            })
    
    def clear_conversation(self):
        if self.current_doc_id and self.current_doc_id in self.conversations:
            self.conversations[self.current_doc_id] = []
    
    def query(self, question):
        if not self.is_ready():
            return "System not ready - check API key status."
        
        current_text = self.get_current_document()
        if not current_text:
            return "Please upload and select a document first."
            
        try:
            with st.spinner("Finding answer..."):
                # Build conversation context
                conversation_history = self.get_conversation_history()
                history_context = ""
                if conversation_history:
                    history_context = "\n\nPrevious questions and answers:\n"
                    for i, conv in enumerate(conversation_history[-3:]):  # Last 3 exchanges
                        history_context += f"Q: {conv['question']}\nA: {conv['answer']}\n\n"
                
                prompt = f"""Based ONLY on the following context:

{current_text}
{history_context}

Question: {question}

Answer based only on the context above:"""
                
                response = self.client.chat.completions.create(
                    model=self.model_name,
                    messages=[{"role": "user", "content": prompt}],
                    max_tokens=500,
                    temperature=0.1
                )
                
                answer = response.choices[0].message.content
                # Store in conversation history
                self.add_to_conversation(question, answer)
                
                return f"**Answer:** {answer}"
        except Exception as e:
            return f"Error: {str(e)}"

# Main App
st.set_page_config(page_title="Cram AI", layout="centered")

st.title("Cram AI")
st.markdown("Upload your study materials and get instant answers")

# Initialize RAG system
if 'rag' not in st.session_state:
    st.session_state.rag = SimpleRAG()

rag = st.session_state.rag

# Initialize conversation counter
if 'conversation_count' not in st.session_state:
    st.session_state.conversation_count = 0

# MULTIPLE FILE UPLOAD
uploaded_files = st.file_uploader(
    "Upload study materials",
    type=['pdf', 'pptx', 'png', 'jpg', 'jpeg'],
    accept_multiple_files=True,
    help="Supported: PDF, PowerPoint, Images"
)

# Process uploaded files
if uploaded_files:
    for uploaded_file in uploaded_files:
        # Check if file already exists
        file_exists = any(doc['filename'] == uploaded_file.name for doc in rag.documents.values())
        
        if not file_exists:
            with st.spinner(f"Processing {uploaded_file.name}..."):
                text = FileProcessor.process_file(uploaded_file)
                
                if text != "Unsupported file type" and len(text.strip()) > 0:
                    doc_id = str(uuid.uuid4())
                    rag.add_document(text, doc_id, uploaded_file.name)
                    st.success(f"✅ {uploaded_file.name} processed")
                else:
                    st.error(f"❌ Could not process {uploaded_file.name}")

# FILE MANAGEMENT DASHBOARD
if rag.is_ready() and rag.documents:
    st.divider()
    st.subheader("Your Documents")
    
    # Document selector
    doc_options = {doc_id: info['filename'] for doc_id, info in rag.documents.items()}
    if doc_options:
        selected_doc = st.selectbox(
            "Select document to query:",
            options=list(doc_options.keys()),
            format_func=lambda x: doc_options[x],
            index=list(doc_options.keys()).index(rag.current_doc_id) if rag.current_doc_id in doc_options else 0
        )
        rag.switch_document(selected_doc)
        
        # Show current document info
        current_doc = rag.documents[selected_doc]
        st.info(f"📄 Currently viewing: {current_doc['filename']}")
        
        col1, col2 = st.columns(2)
        with col1:
            if st.button("🗑️ Clear Conversation", key="clear_conv"):
                rag.clear_conversation()
                st.rerun()
        with col2:
            if st.button("🗑️ Delete Document", key=f"delete_{selected_doc}"):
                rag.delete_document(selected_doc)
                st.rerun()

# CONVERSATION HISTORY - Show previous Q&A
if rag.is_ready() and rag.get_current_document():
    conversation_history = rag.get_conversation_history()
    if conversation_history:
        st.divider()
        st.subheader("Conversation History")
        
        for i, conv in enumerate(conversation_history):
            with st.expander(f"Q: {conv['question'][:50]}...", expanded=(i == len(conversation_history)-1)):
                st.markdown(f"**Question:** {conv['question']}")
                st.markdown(f"**Answer:** {conv['answer']}")

# Q&A SECTION
if rag.is_ready() and rag.get_current_document():
    st.divider()
    st.subheader("Ask Questions")
    
    question = st.text_input("What would you like to know about your document?", key="question_input")
    
    if question:
        st.session_state.conversation_count += 1
        answer = rag.query(question)
        
        # Display the latest answer prominently
        st.markdown("### Latest Answer")
        st.write(answer)
        
        # Auto-refresh to show in conversation history
        st.rerun()
        
elif rag.is_ready() and not rag.documents:
    st.info("📁 Upload some documents to get started!")
elif not rag.is_ready():
    st.warning("⚠️ System initializing...")

# Instructions
with st.expander("How to use Cram AI"):
    st.markdown("""
    1. **Upload** multiple PDFs, PowerPoints, or images
    2. **Select** which document to query
    3. **Ask questions** and see conversation history
    4. **Continue conversations** with follow-up questions
    
    **Supported formats:** PDF, PowerPoint, Images
    """)

# Basic analytics
with st.expander("Developer Info"):
    st.write(f"📊 Documents loaded: {len(rag.documents)}")
    if rag.documents:
        file_types = {}
        for doc in rag.documents.values():
            ext = doc['filename'].split('.')[-1].lower()
            file_types[ext] = file_types.get(ext, 0) + 1
        st.write(f"📁 File types: {file_types}")
    st.write(f"💬 Total conversations: {st.session_state.conversation_count}")

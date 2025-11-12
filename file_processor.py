import os
from PyPDF2 import PdfReader
from pptx import Presentation
import pytesseract
from PIL import Image
import io

class FileProcessor:
    @staticmethod
    def extract_pdf_text(file):
        pdf_reader = PdfReader(file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text
    
    @staticmethod
    def extract_ppt_text(file):
        prs = Presentation(file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    
    @staticmethod
    def extract_image_text(file):
        image = Image.open(file)
        text = pytesseract.image_to_string(image)
        return text
    
    @staticmethod
    def process_file(uploaded_file):
        file_type = uploaded_file.type
        
        if file_type == "application/pdf":
            return FileProcessor.extract_pdf_text(uploaded_file)
        elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return FileProcessor.extract_ppt_text(uploaded_file)
        elif file_type.startswith('image'):
            return FileProcessor.extract_image_text(uploaded_file)
        else:
            return "Unsupported file type"
